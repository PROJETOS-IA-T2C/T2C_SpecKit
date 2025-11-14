"""
Comando init - Cria estrutura inicial do projeto
"""
import os
import shutil
from pathlib import Path
from rich.console import Console
try:
    from importlib.resources import files as resource_files
except ImportError:
    # Python < 3.9 fallback
    from importlib_resources import files as resource_files


def init_project(project_name: str, ai_assistant: str, console: Console):
    """
    Cria a estrutura inicial do projeto RPA Spec-Kit
    
    Args:
        project_name: Nome do projeto
        ai_assistant: AI assistant escolhido (cursor, vscode-copilot, vscode-claude, other)
        console: Console do rich para output
    """
    project_path = Path(project_name)
    
    if project_path.exists():
        raise ValueError(f"Diretório {project_name} já existe!")
    
    console.print(f"[cyan]Criando estrutura do projeto...[/cyan]")
    
    # Criar estrutura de diretórios
    directories = [
        ".specify/memory",
        ".specify/templates",
        ".specify/scripts",
        "specs",
        "generated",
        "DDP",
    ]
    
    # Adicionar diretórios específicos do AI assistant
    if ai_assistant == "cursor":
        directories.append(".cursor/commands")
    elif ai_assistant in ["vscode-copilot", "vscode-claude"]:
        directories.append(".vscode")
        # GitHub Copilot reconhece comandos em .github/prompts/
        if ai_assistant == "vscode-copilot":
            directories.append(".github/prompts")
    
    for directory in directories:
        (project_path / directory).mkdir(parents=True, exist_ok=True)
    
    # Copiar constitution
    console.print("[cyan]Copiando constitution do framework T2C...[/cyan]")
    _copy_constitution(project_path)
    
    # Criar templates vazios
    console.print("[cyan]Criando templates...[/cyan]")
    _create_templates(project_path)
    
    # Criar script de extração de DDP
    console.print("[cyan]Criando script de extração de DDP...[/cyan]")
    _create_extract_ddp_script(project_path)
    
    # Criar requirements.txt
    console.print("[cyan]Criando requirements.txt...[/cyan]")
    _create_requirements_txt(project_path)
    
    # Criar comandos Cursor/VS Code/GitHub Copilot
    if ai_assistant == "cursor":
        console.print("[cyan]Criando comandos Cursor...[/cyan]")
        _create_cursor_commands(project_path)
    elif ai_assistant in ["vscode-copilot", "vscode-claude"]:
        console.print("[cyan]Criando configurações VS Code...[/cyan]")
        _create_vscode_config(project_path, ai_assistant)
        # Criar comandos para GitHub Copilot (reconhece .github/prompts/)
        if ai_assistant == "vscode-copilot":
            console.print("[cyan]Criando comandos GitHub Copilot...[/cyan]")
            _create_github_prompts(project_path)
    
    # Criar arquivos iniciais
    console.print("[cyan]Criando arquivos iniciais...[/cyan]")
    _create_initial_files(project_path, project_name)
    
    console.print("[green]✓[/green] Estrutura criada com sucesso!")


def _copy_constitution(project_path: Path):
    """Copia a constitution do framework T2C do template interno"""
    constitution_path = project_path / ".specify/memory/constitution.md"
    
    try:
        # Usar importlib.resources para acessar arquivos do pacote instalado
        from rpa_speckit import memory
        constitution_resource = resource_files(memory) / "constitution.md"
        
        if constitution_resource.is_file():
            # Ler conteúdo do recurso do pacote
            constitution_content = constitution_resource.read_text(encoding="utf-8")
            constitution_path.write_text(constitution_content, encoding="utf-8")
        else:
            raise FileNotFoundError("Constitution não encontrada no pacote")
    except (ImportError, FileNotFoundError, AttributeError):
        # Fallback: tentar caminho relativo (modo desenvolvimento)
        memory_dir = Path(__file__).parent.parent.parent / "memory"
        internal_constitution = memory_dir / "constitution.md"
        
        if internal_constitution.exists():
            shutil.copy2(internal_constitution, constitution_path)
        else:
            # Se não encontrar, criar versão básica
            basic_constitution = """# Constitution do Framework T2C

Este documento define TODAS as regras, especificações, padrões, exemplos e templates que a IA deve seguir ao gerar código para o framework T2C.

**IMPORTANTE:** Este documento é exclusivamente para uso da IA durante a geração de código.

## Nota

A constitution completa do framework T2C deve estar disponível no SpecKit.

A constitution contém todas as regras, padrões e templates necessários para geração de código.
"""
            constitution_path.write_text(basic_constitution, encoding="utf-8")


def _create_templates(project_path: Path):
    """Cria templates vazios para o desenvolvedor preencher"""
    templates_dir = project_path / ".specify/templates"
    
    # Lista de templates para copiar
    template_files = [
        "spec-template.md",
        "tests-template.md",
        "selectors-template.md",
        "business-rules-template.md",
        "tasks-template.md"
    ]
    
    try:
        # Usar importlib.resources para acessar templates do pacote instalado
        from rpa_speckit import templates
        templates_resource = resource_files(templates)
        
        # Copiar cada template do pacote
        for template_file in template_files:
            dest_template = templates_dir / template_file
            template_resource = templates_resource / template_file
            
            if template_resource.is_file():
                # Ler conteúdo do recurso do pacote
                template_content = template_resource.read_text(encoding="utf-8")
                dest_template.write_text(template_content, encoding="utf-8")
            else:
                # Fallback: criar arquivo vazio se template não existir
                dest_template.write_text(f"# {template_file}\n\n[Template não encontrado no pacote]", encoding="utf-8")
    except (ImportError, AttributeError):
        # Fallback: tentar caminho relativo (modo desenvolvimento)
        internal_templates_dir = Path(__file__).parent.parent.parent / "templates"
        
        for template_file in template_files:
            source_template = internal_templates_dir / template_file
            dest_template = templates_dir / template_file
            
            if source_template.exists():
                shutil.copy2(source_template, dest_template)
            else:
                # Fallback: criar arquivo vazio se template não existir
                dest_template.write_text(f"# {template_file}\n\n[Template não encontrado]", encoding="utf-8")


def _create_extract_ddp_script(project_path: Path):
    """Cria script Python pronto para extração de DDP"""
    scripts_dir = project_path / ".specify/scripts"
    
    # Usar raw string para evitar problemas com escape e encoding
    script_content = r'''#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Script para extração de texto de arquivos DDP.pptx
Este script já está pronto e não deve ser modificado.
"""
import sys
import os
import subprocess
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("python-pptx não está instalado. Instalando automaticamente...", file=sys.stderr)
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx>=0.6.21"], 
                             stdout=sys.stderr, stderr=sys.stderr)
        from pptx import Presentation
        print("python-pptx instalado com sucesso!", file=sys.stderr)
    except Exception as e:
        print(f"Erro ao instalar python-pptx: {e}", file=sys.stderr)
        print("Tente instalar manualmente: pip install python-pptx", file=sys.stderr)
        sys.exit(1)


def extract_ddp(pptx_path: str) -> str:
    """
    Extrai texto de todos os slides de um arquivo DDP.pptx
    
    Args:
        pptx_path: Caminho para o arquivo DDP.pptx (pode ser relativo, absoluto ou apenas nome do arquivo)
        
    Returns:
        Texto formatado com conteúdo de todos os slides
    """
    # Converter para Path e resolver para absoluto (simples e direto)
    pptx_file = Path(pptx_path).resolve()
    
    # Se não encontrar, procurar automaticamente nas pastas comuns
    if not pptx_file.exists():
        # Procurar em DDP/ primeiro
        ddp_dir = Path("DDP")
        if ddp_dir.exists():
            pptx_files = list(ddp_dir.glob("*.pptx"))
            if pptx_files:
                pptx_file = pptx_files[0].resolve()
        
        # Se não encontrou, procurar em specs/*/DDP/
        if not pptx_file.exists():
            for spec_dir in Path("specs").glob("*/DDP"):
                if spec_dir.exists():
                    pptx_files = list(spec_dir.glob("*.pptx"))
                    if pptx_files:
                        pptx_file = pptx_files[0].resolve()
                        break
        
        if not pptx_file.exists():
            raise FileNotFoundError(f"DDP não encontrado: {pptx_path}")
    
    # Usar caminho absoluto sempre (simples)
    presentation = Presentation(str(pptx_file.absolute()))
    
    # Formatar texto para apresentar à LLM
    formatted_text = "# Conteúdo Extraído do DDP\n\n"
    formatted_text += f"**Arquivo:** {pptx_path}\n\n"
    formatted_text += f"**Total de slides:** {len(presentation.slides)}\n\n"
    formatted_text += "---\n\n"
    
    # Passar slide por slide e extrair texto
    for i, slide in enumerate(presentation.slides, 1):
        slide_text = []
        
        # Extrair texto de todas as formas no slide
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        
        # Adicionar slide ao texto formatado
        formatted_text += f"## Slide {i}\n\n"
        formatted_text += "\n".join(slide_text)
        formatted_text += "\n\n---\n\n"
    
    return formatted_text


def main():
    """CLI para extração de DDP"""
    # Configurar encoding UTF-8 para stdout/stderr no Windows
    if sys.platform == 'win32':
        import io
        try:
            sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
            sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')
        except:
            pass
    
    # Se não passou caminho, procurar automaticamente
    if len(sys.argv) < 2:
        # Procurar arquivos .pptx nas pastas comuns
        search_dirs = [Path("DDP")]
        for spec_dir in Path("specs").glob("*/DDP"):
            search_dirs.append(spec_dir)
        
        pptx_file = None
        for search_dir in search_dirs:
            if search_dir.exists():
                pptx_files = list(search_dir.glob("*.pptx"))
                if pptx_files:
                    pptx_file = pptx_files[0].resolve()
                    break
        
        if not pptx_file:
            print("Erro: Nenhum arquivo .pptx encontrado. Use: python .specify/scripts/extract-ddp.py <caminho>", file=sys.stderr)
            sys.exit(1)
        
        ddp_path = str(pptx_file)
    else:
        ddp_path = sys.argv[1]
    
    try:
        extracted_text = extract_ddp(ddp_path)
        print(extracted_text)
    except FileNotFoundError as e:
        print(f"Erro: {e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"Erro ao extrair DDP: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
'''
    
    extract_script = scripts_dir / "extract-ddp.py"
    extract_script.write_text(script_content, encoding="utf-8")


def _create_requirements_txt(project_path: Path):
    """Cria requirements.txt com dependências necessárias"""
    requirements_content = """# Dependências para scripts do projeto
python-pptx>=0.6.21
"""
    (project_path / "requirements.txt").write_text(requirements_content, encoding="utf-8")


def _get_command_content(command_name: str) -> str:
    """Retorna o conteúdo completo de um comando (reutilizável para Cursor e VS Code)"""
    commands = {
        "t2c.extract-ddp": """# Extrair DDP

Extrai o texto de todos os slides de um arquivo DDP.pptx para que a LLM possa preencher os arquivos de especificação.

## Uso

\`\`\`
/t2c.extract-ddp [caminho_do_ddp]
\`\`\`

## Exemplo

\`\`\`
/t2c.extract-ddp specs/001-automacao-exemplo/DDP/ddp.pptx
\`\`\`

## 🚨 REGRA FUNDAMENTAL - LEITURA CUIDADOSA DO DDP

**⚠️ EXTREMAMENTE CRÍTICO - SEM ISSO TUDO ESTARÁ ERRADO:**

**ANTES DE QUALQUER OUTRA AÇÃO, a LLM DEVE:**

1. **⚠️ OBRIGATÓRIO: Ler o DDP com ATENÇÃO TOTAL** (localizado em `specs/001-[nome]/DDP/ddp.pptx` ou caminho fornecido)
   - Ler o DDP **COMPLETO** do início ao fim, **palavra por palavra**
   - **NÃO pular NENHUMA seção** - mesmo que pareça irrelevante
   - **NÃO fazer suposições** - se algo não está claro, revisar o DDP
   - Ler **múltiplas vezes** se necessário para garantir compreensão completa
   - **Identificar TUDO** que está mapeado no DDP

2. **⚠️ OBRIGATÓRIO: Mapear COMPLETAMENTE o DDP:**
   - Criar uma lista de **TODAS as etapas** identificadas (INIT, FILA, LOOP STATION, END PROCESS)
   - Criar uma lista de **TODAS as exceções de negócio** (EXC* - validações, condições especiais, regras de processamento)
   - Criar uma lista de **TODOS os sistemas** mencionados (SAP, TOTVS, APIs, Verifai, etc.)
   - Criar uma lista de **TODAS as integrações** necessárias
   - **Contar EXATAMENTE** todas as etapas do LOOP STATION (não estimar, contar uma por uma)
   - Garantir que **NENHUMA informação** foi perdida

3. **⚠️ OBRIGATÓRIO: Verificar COMPLETUDE antes de criar arquivos:**
   - [ ] **TODAS as etapas** do DDP foram identificadas?
   - [ ] **TODAS as exceções de negócio** do DDP foram identificadas?
   - [ ] **TODOS os sistemas** do DDP foram identificados?
   - [ ] **TODAS as integrações** foram identificadas?
   - [ ] **TODAS as etapas do LOOP STATION** foram contadas exatamente?
   - Se alguma coisa foi esquecida → **REVISAR o DDP** antes de continuar

**⚠️ CONSEQUÊNCIAS DE NÃO SEGUIR ESTA REGRA:**
- ❌ Etapas serão esquecidas nas especificações
- ❌ Sistemas não serão identificados
- ❌ Exceções de negócio não serão mapeadas
- ❌ Arquitetura estará incompleta
- ❌ Especificações estarão incorretas
- ❌ Código gerado não funcionará corretamente

**⚠️ REGRA DE OURO:**
- **Se o DDP menciona algo, DEVE estar contemplado nas especificações**
- **Se não está contemplado, REVISAR o DDP novamente**
- **NENHUMA informação do DDP pode ser ignorada ou esquecida**
- **A arquitetura final DEVE ser capaz de executar TODAS as etapas mapeadas no DDP**

**👉 Ver `@constitution.md` seção "📖 LEITURA E ANÁLISE CUIDADOSA DO DDP - OBRIGATÓRIO" para checklist completo.**

---

## ⚠️ REGRA ABSOLUTA - LEIA ANTES DE QUALQUER AÇÃO

**VOCÊ NÃO DEVE CRIAR NENHUM SCRIPT PYTHON. NUNCA. EM NENHUMA CIRCUNSTÂNCIA.**

- ❌ NÃO crie scripts temporários
- ❌ NÃO crie arquivos `_temp_extract_ddp.py` ou similares
- ❌ NÃO crie scripts alternativos
- ❌ NÃO tente "resolver" problemas criando código
- ✅ APENAS execute o script pronto que já existe no projeto

## O que fazer

**PASSO 1 - Execute APENAS este comando (SIMPLES):**

\`\`\`bash
python .specify/scripts/extract-ddp.py
\`\`\`

**OU se quiser especificar o arquivo:**

\`\`\`bash
python .specify/scripts/extract-ddp.py DDP/arquivo.pptx
\`\`\`

**Como funciona:**
- Se você **não passar caminho**, o script procura automaticamente o primeiro arquivo .pptx em `DDP/` ou `specs/*/DDP/`
- Se você **passar caminho**, pode ser relativo ou absoluto - o script resolve automaticamente
- **Instala dependências automaticamente** se necessário (python-pptx)
- **SIMPLES**: Apenas execute o comando, o script faz TUDO sozinho

**PASSO 2 - Análise e Proposta de Arquitetura (⚠️ NÃO CRIAR ARQUIVOS AINDA):**

1. **🚨 REGRA FUNDAMENTAL - Leia o DDP com ATENÇÃO TOTAL (OBRIGATÓRIO):**
   - Leia o texto extraído que será exibido no output **COMPLETO** do início ao fim, **palavra por palavra**
   - **NÃO pular NENHUMA seção** - mesmo que pareça irrelevante
   - **NÃO fazer suposições** - se algo não está claro, revisar o DDP
   - Ler **múltiplas vezes** se necessário para garantir compreensão completa
   - **NÃO DEIXE PASSAR NENHUMA ETAPA, REGRA, SISTEMA OU EXCEÇÃO** mapeada no DDP
   - Identifique **TODAS as etapas** (INIT, FILA, LOOP STATION, END PROCESS)
   - Identifique **TODAS as exceções de negócio** (EXC* - validações, condições especiais, regras de processamento - tudo que pode gerar uma exceção ou regra específica)
   - Identifique **TODOS os sistemas** (APIs, UI, Verifai, etc.)
   - Identifique **TODAS as integrações** necessárias
   - **Conte EXATAMENTE** todas as etapas do LOOP STATION (não estime, conte uma por uma)
   - Crie uma lista **ESCRITA** de **TODAS as informações** identificadas (etapas, exceções, sistemas, integrações)
   - Garantir que **NENHUMA informação** foi perdida

2. **⚠️ OBRIGATÓRIO: Consulte o `@constitution.md`** (localizado em `src/rpa_speckit/memory/constitution.md`) para decidir a arquitetura
   - **PRIMEIRO:** Leia a **seção "🚨 REGRA FUNDAMENTAL - LEITURA CUIDADOSA DO DDP"** no início do documento - Esta é EXTREMAMENTE CRÍTICA
   - **SEGUNDO:** Leia a **seção 0: 🚨 REGRA CRÍTICA - SEGUIR ESTRUTURA DOS TEMPLATES EXATAMENTE** - Esta é EXTREMAMENTE IMPORTANTE
   - **TERCEIRO:** Leia especialmente a **PARTE 1.5: Arquitetura de Robôs** (ou seção 13) para decisão de arquitetura
   - Leia a seção **"📖 LEITURA E ANÁLISE CUIDADOSA DO DDP - OBRIGATÓRIO"** e siga o checklist obrigatório COMPLETO
   - Verifique as **REGRAS OBRIGATÓRIAS DE SEPARAÇÃO** primeiro
   - Se QUALQUER regra obrigatória se aplicar → SEPARAR É OBRIGATÓRIO
   - **⚠️ IMPORTANTE:** Quanto mais quebrar os robôs seguindo as regras, melhor ainda. Seja proativo em separar quando as regras se aplicam.

3. **📋 CRIAR PROPOSTA DE ARQUITETURA (⚠️ NÃO CRIAR ARQUIVOS AINDA - APENAS APRESENTAR PROPOSTA):**
   
   **A LLM DEVE criar um resumo compacto mas detalhado apresentando ao usuário:**
   
   **a) Resumo do Processo:**
   - **Sistemas envolvidos:** Lista de todos os sistemas (APIs, UI, Verifai, Excel, Word, portais, etc.)
   - **Fluxo geral:** Descrição breve do processo do início ao fim
   - **Complexidade identificada:** Pontos de complexidade (múltiplas APIs, LOOP extenso, sistemas diferentes, etc.)
   - **Exceções de negócio:** Quantidade e tipos de exceções identificadas
   
   **b) Decisão de Arquitetura:**
   - **Quantidade de robôs propostos:** X robôs (1, 2, 3, 4, 5 ou quantos forem necessários)
   - **Tipo de cada robô:** Standalone / Dispatcher / Performer
   - **Justificativa detalhada:**
     - Quais regras obrigatórias se aplicaram (REGRA 1, REGRA 2, REGRA 3, REGRA 4, REGRA 5)
     - Por que decidiu separar dessa forma específica
     - Benefícios da separação (isolamento de erros, execução retroativa, manutenibilidade, etc.)
     - **⚠️ IMPORTANTE:** Se houver dúvida entre separar mais ou menos, SEMPRE optar por separar mais (seguindo as regras). Quanto mais quebrar os robôs seguindo as regras, melhor ainda.
   
   **c) Detalhamento por Robô:**
   
   Para cada robô proposto, mostrar de forma clara e organizada:
   - **Nome/Role:** O que este robô faz (ex: "Dispatcher - Prepara dados do Pipefy, consulta APIs e envia para Verifai")
   - **Etapas principais:** Lista das etapas principais que este robô executará (INIT, FILA, LOOP STATION, END PROCESS)
   - **Sistemas que utiliza:** Quais sistemas este robô interage (ex: "Pipefy API, CNPJ API, Verifai")
   - **Entrada:** De onde recebe dados (fila própria, fila compartilhada, Excel, API, etc.)
   - **Saída:** O que produz (popula fila do próximo robô, finaliza processo, etc.)
   - **Ordem na cadeia:** Posição na sequência de execução (1, 2, 3, etc.)
   
   **d) Estrutura de Pastas:**
   - Mostrar a estrutura de pastas que será criada (robot1/, robot2/, etc. ou raiz se standalone)
   
   **Formato sugerido para apresentação ao usuário:**
   ```markdown
   ## 📋 Proposta de Arquitetura
   
   ### 📊 Resumo do Processo
   - **Sistemas envolvidos:** [lista de sistemas]
   - **Fluxo geral:** [descrição breve]
   - **Complexidade:** [pontos de complexidade identificados]
   - **Exceções de negócio:** [quantidade e tipos]
   
   ### 🏗️ Decisão de Arquitetura
   - **Quantidade de robôs:** X robôs
   - **Justificativa:**
     - [Regra obrigatória aplicada]: [explicação]
     - [Regra obrigatória aplicada]: [explicação]
     - Benefícios: [isolamento de erros, execução retroativa, etc.]
   
   ### 🤖 Detalhamento por Robô
   
   **Robot1 (Dispatcher):**
   - **Função:** [descrição clara do papel]
   - **Etapas principais:**
     - INIT: [o que faz no INIT]
     - FILA: [o que faz na FILA]
     - LOOP STATION: [o que faz no LOOP]
     - END PROCESS: [o que faz no final]
   - **Sistemas:** [lista de sistemas]
   - **Entrada:** [fonte de dados]
   - **Saída:** [o que produz/popula]
   - **Ordem:** 1
   
   **Robot2 (Performer):**
   [mesmo formato]
   
   ### 📂 Estrutura Proposta
   ```
   specs/001-[nome]/
   ├── robot1/
   │   ├── spec.md
   │   ├── selectors.md
   │   ├── business-rules.md
   │   └── tests.md
   ├── robot2/
   │   ├── spec.md
   │   └── ...
   └── tasks.md (será criado com /t2c.tasks)
   ```
   
   **⏳ Aguardando sua aprovação para criar os arquivos...**
   ```

4. **✅ APROVAÇÃO DO USUÁRIO:**
   
   **Após apresentar a proposta, a LLM DEVE:**
   - **Aguardar explicitamente** a aprovação do usuário
   - **NÃO criar arquivos** até receber aprovação
   - Se o usuário der **OK, aprovar, confirmar** ou similar → Prosseguir para PASSO 3 (criar arquivos)
   - Se o usuário **NÃO aprovar** ou pedir ajustes:
     - Entender o feedback do usuário
     - Ajustar a proposta conforme solicitado
     - Apresentar nova proposta
     - Repetir até aprovação
   - Se o usuário questionar a decisão, explicar melhor a justificativa baseada nas regras do `constitution.md`

**PASSO 3 - Criação dos Arquivos (⚠️ APENAS APÓS APROVAÇÃO DO USUÁRIO):**

1. **🚨 VERIFICAÇÃO FINAL OBRIGATÓRIA - ANTES DE CRIAR OS ARQUIVOS:**
   - [ ] **TODAS as etapas** do DDP foram identificadas e estão contempladas na arquitetura aprovada?
   - [ ] **TODAS as exceções de negócio** (EXC* - validações, condições especiais, regras de processamento) do DDP foram identificadas e estarão mapeadas no business-rules.md?
   - [ ] **TODOS os sistemas** mencionados no DDP foram identificados e estão contemplados?
   - [ ] **TODAS as integrações** necessárias foram identificadas e estão consideradas?
   - [ ] **TODAS as exceções** mapeadas no DDP foram identificadas e estão contempladas?
   - [ ] **TODAS as etapas do LOOP STATION** foram contadas EXATAMENTE (não estimadas) e estarão no spec.md?
   - [ ] **NENHUMA informação** do DDP foi esquecida ou ignorada?
   - [ ] A arquitetura aprovada contempla **TODAS as regras obrigatórias** que se aplicam?
   - Se **QUALQUER coisa** do DDP não foi contemplada → **REVISAR o DDP COMPLETO novamente** e **CORRIGIR** antes de criar os arquivos
   - **REGRA DE OURO:** Se o DDP menciona algo, DEVE estar contemplado nas especificações

2. **Criar os arquivos de especificação** baseado na arquitetura aprovada pelo usuário:
   - Seguir EXATAMENTE a estrutura dos templates (ver PASSO 3 abaixo)
   - Preencher cada arquivo baseado no conteúdo do DDP
   - Garantir que **TODAS as etapas, TODAS as exceções, TODOS os sistemas e TODAS as regras** estejam contempladas

## Arquivos a preencher

**🚨 REGRA CRÍTICA - NÃO CRIAR tasks.md:**
- ❌ **NÃO criar** \`tasks.md\` neste comando
- ❌ **NÃO gerar** tasks.md automaticamente
- ✅ **tasks.md** deve ser criado APENAS quando o usuário executar o comando `/t2c.tasks`
- ✅ **Aguardar** o comando explícito do usuário para gerar tasks.md

### Se Standalone (1 robô):
- \`specs/001-[nome]/spec.md\` - Especificação técnica e arquitetura (ARQUIVO PRINCIPAL)
- \`specs/001-[nome]/tests.md\` - Cenários de teste e validações
- \`specs/001-[nome]/selectors.md\` - Seletores Clicknium
- \`specs/001-[nome]/business-rules.md\` - Regras de negócio
- ❌ **NÃO criar** \`tasks.md\` - será criado apenas com o comando `/t2c.tasks`

### Se Múltiplos Robôs (quando regra obrigatória se aplicar):
- \`specs/001-[nome]/robot1/spec.md\` - Especificação do robô 1 (Dispatcher)
- \`specs/001-[nome]/robot1/tests.md\` - Testes do robô 1
- \`specs/001-[nome]/robot1/selectors.md\` - Seletores do robô 1
- \`specs/001-[nome]/robot1/business-rules.md\` - Regras de negócio do robô 1
- \`specs/001-[nome]/robot2/spec.md\` - Especificação do robô 2 (Performer)
- \`specs/001-[nome]/robot2/tests.md\` - Testes do robô 2
- \`specs/001-[nome]/robot2/selectors.md\` - Seletores do robô 2
- \`specs/001-[nome]/robot2/business-rules.md\` - Regras de negócio do robô 2
- ❌ **NÃO criar** \`tasks.md\` - será criado apenas com o comando `/t2c.tasks`

**⚠️ IMPORTANTE:** 
- Se houver múltiplos robôs, **NÃO criar** \`spec.md\` na raiz
- Cada robô tem seu próprio \`spec.md\` dentro de sua pasta (\`robot1/\`, \`robot2/\`)
- **NUNCA criar tasks.md** neste comando - aguardar comando `/t2c.tasks` do usuário

## Detalhes dos arquivos

- **spec.md**: ARQUIVO PRINCIPAL - Definir arquitetura completa (INIT, FILA, LOOP STATION, END PROCESS), stack tecnológica, integrações, estrutura de dados. **DEVE incluir seção "Arquitetura de Robôs" no início** conforme `constitution.md` seção 13.
- **tests.md**: Extrair cenários de usuário, requisitos funcionais/não-funcionais, critérios de sucesso, entidades principais
- **selectors.md**: Identificar elementos de UI mencionados no DDP (botões, campos, tabelas, etc.)
- **business-rules.md**: Extrair todas as exceções de negócio (EXC* - validações, condições especiais, regras de processamento - tudo que pode gerar uma exceção ou regra específica)

## 🚨 REGRA CRÍTICA - SEGUIR ESTRUTURA DO TEMPLATE EXATAMENTE

**⚠️ EXTREMAMENTE IMPORTANTE - OBRIGATÓRIO:**

Ao criar ou atualizar qualquer arquivo de especificação (especialmente `spec.md`), a LLM DEVE:

1. **⚠️ OBRIGATÓRIO: Ler o template correspondente ANTES de criar o arquivo:**
   - Para `spec.md`: Ler `@spec-template.md` (localizado em `.specify/templates/spec-template.md`)
   - Para `tests.md`: Ler `@tests-template.md`
   - Para `selectors.md`: Ler `@selectors-template.md`
   - Para `business-rules.md`: Ler `@business-rules-template.md`

2. **⚠️ OBRIGATÓRIO: Replicar a estrutura EXATAMENTE como está no template:**
   - **TODAS as seções** do template devem estar presentes no arquivo criado
   - **TODOS os títulos** (headers) devem estar presentes na mesma ordem
   - **TODOS os subtítulos** devem estar presentes
   - **TODOS os campos** (campos com `- **Campo:**`) devem estar presentes
   - **A ordem das seções** deve ser EXATAMENTE a mesma do template
   - **A formatação** (markdown, listas, blocos de código) deve ser EXATAMENTE a mesma

3. **⚠️ OBRIGATÓRIO: Remover apenas anotações e exemplos:**
   - Remover textos entre colchetes `[...]`
   - Remover exemplos que não são do processo real
   - Remover anotações explicativas
   - **MANTER** toda a estrutura, seções, títulos e campos do template

4. **⚠️ OBRIGATÓRIO: Preencher com informações reais do DDP:**
   - Substituir `[Nome da Etapa]` por nomes reais das etapas do DDP
   - Substituir `[Descrição]` por descrições reais do DDP
   - Substituir `[N/A]` apenas quando realmente não se aplica
   - Preencher TODOS os campos com informações reais do processo

5. **⚠️ VERIFICAÇÃO OBRIGATÓRIA antes de finalizar:**
   - [ ] Todas as seções do template estão presentes?
   - [ ] Todos os títulos estão na mesma ordem do template?
   - [ ] Todos os campos obrigatórios estão preenchidos?
   - [ ] A estrutura está EXATAMENTE igual ao template (apenas sem anotações/exemplos)?
   - [ ] Não foi adicionada nenhuma seção que não está no template?
   - [ ] Não foi removida nenhuma seção que está no template?

**Exemplo de estrutura correta para spec.md:**
```
1. Título principal (igual ao template)
2. Aviso sobre remover anotações (igual ao template)
3. Seção "Arquitetura de Robôs" (igual ao template)
4. Seção "Stack Tecnológica" (igual ao template)
5. Seção "Visão Geral do Fluxo" (igual ao template)
6. Seção "INIT: Inicialização" (igual ao template)
7. Seção "FILA: Preenchimento da Fila" (igual ao template)
8. Seção "LOOP STATION: Processamento Principal" (igual ao template)
9. Seção "END PROCESS: Finalização" (igual ao template)
10. Seção "Integrações do Projeto" (igual ao template)
11. Seção "Observações Gerais" (igual ao template)
```

**❌ NÃO FAZER:**
- ❌ Criar seções que não estão no template
- ❌ Remover seções que estão no template
- ❌ Alterar a ordem das seções
- ❌ Modificar a estrutura de campos
- ❌ Adicionar campos que não estão no template
- ❌ Remover campos que estão no template

**✅ FAZER:**
- ✅ Seguir EXATAMENTE a estrutura do template
- ✅ Manter TODAS as seções do template
- ✅ Preencher com informações reais do DDP
- ✅ Remover apenas anotações e exemplos (textos entre `[...]`)

## ⚠️ REGRAS CRÍTICAS

**🚨 REGRA FUNDAMENTAL - Leitura do DDP:**
- **SEMPRE leia o DDP com ATENÇÃO TOTAL** - palavra por palavra, do início ao fim
- **NÃO DEIXE PASSAR NENHUMA ETAPA, REGRA, SISTEMA OU EXCEÇÃO**
- Se algo do DDP não for contemplado, **TODAS as especificações estarão incorretas**

**🚨 REGRA CRÍTICA - Separação de Robôs:**
- **Quanto mais quebrar os robôs seguindo as regras, melhor ainda**
- Seja **proativo** em separar quando as regras se aplicam
- Se houver dúvida entre separar mais ou menos, **SEMPRE optar por separar mais** (seguindo as regras)
- Não seja conservador - se as regras obrigatórias se aplicam, **SEPARAR É OBRIGATÓRIO**
- Após aplicar REGRA 5 (Verifai), **SEMPRE verificar** se REGRA 1 também se aplica ao processamento subsequente

**🚨 REGRA CRÍTICA - Aprovação Antes de Criar:**
- **NUNCA criar arquivos** antes de apresentar a proposta e receber aprovação do usuário
- **SEMPRE apresentar** a proposta de arquitetura primeiro
- **AGUARDAR** explicitamente a aprovação do usuário
- Se o usuário pedir ajustes, ajustar a proposta e apresentar novamente

**🚨 REGRA CRÍTICA - Estrutura dos Templates:**
- **SEMPRE consulte o template correspondente** (ex: `@spec-template.md`) ANTES de criar qualquer arquivo
- **SEMPRE replique a estrutura EXATAMENTE** como está no template
- **SEMPRE consulte o `@constitution.md`** seção 0 sobre seguir estrutura dos templates

**🚨 REGRA ABSOLUTA - tasks.md:**
- ❌ **NUNCA criar** \`tasks.md\` neste comando
- ❌ **NÃO gerar** tasks.md automaticamente
- ✅ **tasks.md** será criado APENAS quando o usuário executar explicitamente o comando `/t2c.tasks`
- ✅ **Aguardar** o comando do usuário - não antecipar a criação de tasks.md

## Lembre-se

- O script `.specify/scripts/extract-ddp.py` JÁ EXISTE no projeto e está pronto - apenas execute-o
- Use os templates em \`.specify/templates/\` como referência para a estrutura
- Mantenha a numeração das regras (EXC001, EXC002, etc.)
- Se os arquivos já existirem, atualize-os com as novas informações do DDP, mas **MANTENHA a estrutura do template**
- **SEMPRE verifique** que **TODAS as etapas, TODAS as exceções, TODOS os sistemas e TODAS as regras** do DDP estão contempladas antes de criar os arquivos""",
        "t2c.tasks": """# Gerar Tasks

Gera o arquivo tasks.md baseado em spec.md e business-rules.md, incluindo estimativas de tempo para cada tarefa.

## Uso

\`\`\`
/t2c.tasks [caminho_da_spec]
\`\`\`

## Exemplo

\`\`\`
/t2c.tasks specs/001-automacao-exemplo
\`\`\`

## O que faz

1. Lê spec.md e business-rules.md
2. Analisa os requisitos e regras
3. Gera breakdown de tarefas organizado por fases:
   - Init (T2CInitAllApplications)
   - Process (T2CProcess)
   - End Process (T2CCloseAllApplications)
4. **Calcula estimativas de tempo** para cada tarefa (considerando desenvolvedor pleno)
5. Cria tasks.md com:
   - Tabela de visão geral de estimativas no início
   - Cada tarefa com sua estimativa de tempo e justificativa
   - Resumo executivo com métricas de tempo

## Arquivo Gerado

- \`specs/001-[nome]/tasks.md\` com:
  - Tabela de visão geral (resumo executivo, top 5 tasks, estimativas por fase/robô)
  - Tasks detalhadas com estimativas individuais

## Estimativas de Tempo

**⚠️ OBRIGATÓRIO - Consultar Base de Dados de Complexidade:**

Antes de calcular qualquer estimativa, a LLM DEVE:

1. **Consultar o arquivo `@system_complexity.json`** (localizado em `src/rpa_speckit/memory/system_complexity.json`)
   - Este arquivo contém multiplicadores objetivos baseados em dados reais
   - NÃO fazer estimativas baseadas em suposições - sempre consultar a base de dados

2. **Identificar os sistemas mencionados no spec.md:**
   - Verificar se o sistema está listado na base de dados (sistemas conhecidos)
   - Se não estiver, classificar por categoria (portal governo, legado, menos conhecido, customizado)

3. **Aplicar multiplicadores conforme a base de dados:**
   - Multiplicador do sistema (baseado na categoria ou sistema específico)
   - Multiplicador de interface (Web Moderna, Web Legado, Desktop, Terminal)
   - Multiplicador de documentação (Completa, Parcial, Sem documentação)
   - Multiplicador de seletores (Estáveis, Instáveis, Dinâmicos)

4. **Calcular estimativa final:**
   ```
   Estimativa Final = Estimativa Base × Multiplicador Sistema × Multiplicador Interface × Multiplicador Documentação × Multiplicador Seletores
   ```

5. **Documentar na justificativa:**
   - Sempre mencionar os multiplicadores aplicados da base de dados
   - Explicar por que cada multiplicador foi usado
   - Referenciar o sistema e categoria aplicada

**Regras de Estimativa:**
- **Base:** Desenvolvedor pleno (não mencionar isso no documento, apenas usar como referência)
- **Formato:** Horas (ex: "2 horas", "4 horas", "0.5 horas")
- **Justificativa:** DEVE incluir referência aos multiplicadores aplicados da base de dados
- **Tabela de visão geral:** Inclui tempo total, top 5 tasks mais demoradas, distribuição por fase e por robô

**⚠️ IMPORTANTE:** 
- NUNCA fazer estimativas sem consultar `@system_complexity.json`
- SEMPRE documentar quais multiplicadores foram aplicados
- Ver seção 14 do `@constitution.md` para instruções detalhadas sobre como usar a base de dados

## Notas

- Este comando é opcional - o desenvolvedor pode criar tasks.md manualmente
- As tarefas geradas devem ser revisadas e ajustadas conforme necessário
- As estimativas são baseadas na complexidade descrita no spec.md e business-rules.md""",
        "t2c.implement": """# Implementar Framework T2C

Gera o framework T2C completo baseado nas especificações preenchidas.

## Uso

\`\`\`
/t2c.implement [caminho_da_spec] [--robot nome_do_robo]
\`\`\`

## Exemplos

\`\`\`
# Gerar todos os robôs (ou standalone)
/t2c.implement specs/001-automacao-exemplo

# Gerar apenas um robô específico (se múltiplos robôs)
/t2c.implement specs/001-automacao-exemplo --robot robot1
/t2c.implement specs/001-automacao-exemplo --robot robot2
\`\`\`

## Estrutura de Robôs

O comando detecta automaticamente se o projeto é:
- **Standalone**: Um único robô (spec.md na raiz)
- **Múltiplos robôs**: Vários robôs (robot1/, robot2/, etc.)

### Standalone
\`\`\`
specs/001-[nome]/
├── spec.md
├── selectors.md
├── business-rules.md
├── tests.md
└── tasks.md
\`\`\`

### Múltiplos Robôs
\`\`\`
specs/001-[nome]/
├── robot1/
│   ├── spec.md
│   ├── selectors.md
│   ├── business-rules.md
│   └── tests.md
├── robot2/
│   ├── spec.md
│   ├── selectors.md
│   ├── business-rules.md
│   └── tests.md
└── tasks.md  # Compartilhado
\`\`\`

## O que faz

1. Detecta estrutura (standalone ou múltiplos robôs)
2. Valida se todos os arquivos necessários estão preenchidos:
   - spec.md (ARQUIVO PRINCIPAL - arquitetura completa)
   - selectors.md
   - business-rules.md
   - tests.md
   - tasks.md (compartilhado se múltiplos robôs)
   - config/*.md
3. Baixa o framework T2C do GitHub (organização privada)
4. Gera estrutura completa:
   - Standalone: \`generated/[nome-automacao]/\`
   - Múltiplos: \`generated/[nome-automacao]-robot1/\`, \`generated/[nome-automacao]-robot2/\`, etc.
5. Copia arquivos do framework base
6. Gera arquivos customizados para cada robô:
   - bot.py
   - T2CProcess.py
   - T2CInitAllApplications.py
   - T2CCloseAllApplications.py
   - Config.xlsx
7. Substitui variáveis de template
8. Gera requirements.txt, setup.py, README.md para cada robô

## Parâmetros

- \`caminho_da_spec\`: Caminho para o diretório da spec (ex: specs/001-automacao-exemplo)
- \`--robot nome_do_robo\`: (Opcional) Gera apenas o robô especificado (ex: robot1, robot2). Se não especificado, gera todos os robôs.

## Arquivos Gerados

- **Standalone**: Estrutura completa em \`generated/[nome-automacao]/\`
- **Múltiplos**: Estrutura completa em \`generated/[nome-automacao]-robot1/\`, \`generated/[nome-automacao]-robot2/\`, etc.

## Pré-requisitos

- Acesso ao repositório privado do framework T2C
- Git configurado
- Python 3.8+ instalado

## Notas

- O framework é gerado do zero a cada execução
- Arquivos customizados são gerados baseados nas specs de cada robô
- Arquivos do framework base são copiados (não modificados)
- Se múltiplos robôs, cada um tem seu próprio framework completo gerado""",
        "t2c.validate": """# Validar Especificações

Valida a estrutura e completude dos arquivos de especificação.

## Uso

\`\`\`
/t2c.validate [caminho_da_spec]
\`\`\`

## Exemplo

\`\`\`
/t2c.validate specs/001-automacao-exemplo
\`\`\`

## O que faz

1. Verifica se todos os arquivos necessários existem:
   - spec.md (ARQUIVO PRINCIPAL)
   - selectors.md
   - business-rules.md
   - tasks.md
   - config/*.md
2. Valida estrutura dos arquivos
3. Verifica se campos obrigatórios estão preenchidos
4. Gera relatório de validação

## Saída

Relatório indicando:
- ✓ Arquivos presentes
- ✓ Campos preenchidos
- ✗ Arquivos faltando
- ✗ Campos obrigatórios vazios

## Notas

- Execute antes de /t2c.implement para garantir que tudo está pronto
- Corrija os problemas indicados antes de prosseguir"""
    }
    return commands.get(command_name, "")


def _create_cursor_commands(project_path: Path):
    """Cria comandos Cursor"""
    commands_dir = project_path / ".cursor/commands"
    
    # Usar a mesma função para garantir conteúdo idêntico
    for cmd_name in ["t2c.extract-ddp", "t2c.tasks", "t2c.implement", "t2c.validate"]:
        content = _get_command_content(cmd_name)
        (commands_dir / f"{cmd_name}.md").write_text(content, encoding="utf-8")


def _create_github_prompts(project_path: Path):
    """Cria comandos para GitHub Copilot usando .github/prompts/ com extensão .prompt.md"""
    prompts_dir = project_path / ".github" / "prompts"
    
    # GitHub Copilot requer extensão .prompt.md (não apenas .md)
    for cmd_name in ["t2c.extract-ddp", "t2c.tasks", "t2c.implement", "t2c.validate"]:
        content = _get_command_content(cmd_name)
        # Copilot reconhece arquivos .prompt.md em .github/prompts/
        (prompts_dir / f"{cmd_name}.prompt.md").write_text(content, encoding="utf-8")


def _create_vscode_config(project_path: Path, ai_assistant: str):
    """Cria configurações VS Code - apenas settings.json para GitHub Copilot"""
    vscode_dir = project_path / ".vscode"
    
    # Criar settings.json
    settings = {
        "python.defaultInterpreterPath": "${workspaceFolder}/.venv/bin/python",
        "python.analysis.typeCheckingMode": "basic",
        "files.exclude": {
            "**/__pycache__": True,
            "**/*.pyc": True,
        }
    }
    
    if ai_assistant == "vscode-copilot":
        settings["github.copilot.enable"] = {
            "*": True
        }
        # Configurações para Copilot Chat reconhecer slash commands
        settings["github.copilot.chat.enable"] = True
        
        # Git autofetch para manter contexto atualizado
        settings["git.autofetch"] = True
        
        # Mapear comandos para arquivos .prompt.md em .github/prompts/
        # Isso permite autocomplete e reconhecimento automático dos slash commands
        settings["chat.promptFilesRecommendations"] = {
            "t2c.extract-ddp": True,
            "t2c.tasks": True,
            "t2c.implement": True,
            "t2c.validate": True
        }
        
        # Permitir execução automática de scripts em .specify/scripts/
        # Isso evita pedir confirmação a cada execução de script
        settings["chat.tools.terminal.autoApprove"] = {
            ".specify/scripts/": True,
            ".specify/scripts/bash/": True,
            ".specify/scripts/powershell/": True
        }
    
    import json
    (vscode_dir / "settings.json").write_text(
        json.dumps(settings, indent=2),
        encoding="utf-8"
    )
    
    # Criar tasks.json (para executar scripts via VS Code Tasks como alternativa)
    _create_vscode_tasks(vscode_dir)


def _create_vscode_commands(commands_dir: Path):
    """Cria arquivos markdown de comandos EXATAMENTE como no Cursor (com slash commands)"""
    
    # Usar a mesma função para garantir conteúdo idêntico
    for cmd_name in ["t2c.extract-ddp", "t2c.tasks", "t2c.implement", "t2c.validate"]:
        content = _get_command_content(cmd_name)
        (commands_dir / f"{cmd_name}.md").write_text(content, encoding="utf-8")


def _create_copilot_instructions(vscode_dir: Path, commands_dir: Path):
    """Cria arquivo de instruções do Copilot dentro de .vscode para suportar slash commands"""
    instructions_content = """# GitHub Copilot Instructions - T2C Commands

Este projeto usa comandos slash customizados (similar ao Cursor) que devem ser reconhecidos pelo GitHub Copilot Chat.

## Comandos Disponíveis

Quando o usuário digitar um comando slash no chat do Copilot, você deve:

1. **Reconhecer o comando**: Se o usuário digitar `/t2c.extract-ddp`, `/t2c.tasks`, `/t2c.implement`, ou `/t2c.validate`
2. **Ler o arquivo correspondente**: Consulte `.vscode/commands/[nome-do-comando].md` para entender o que fazer
3. **Executar as instruções**: Siga EXATAMENTE as instruções do arquivo markdown

## Comandos Slash Customizados

### `/t2c.extract-ddp [caminho]`
- **Arquivo de referência**: `.vscode/commands/t2c.extract-ddp.md`
- **Função**: Extrai texto de arquivos DDP.pptx
- **Uso**: `/t2c.extract-ddp` ou `/t2c.extract-ddp specs/001-exemplo/DDP/ddp.pptx`

### `/t2c.tasks [caminho]`
- **Arquivo de referência**: `.vscode/commands/t2c.tasks.md`
- **Função**: Gera arquivo tasks.md baseado em spec.md e business-rules.md
- **Uso**: `/t2c.tasks specs/001-exemplo`

### `/t2c.implement [caminho]`
- **Arquivo de referência**: `.vscode/commands/t2c.implement.md`
- **Função**: Gera framework T2C completo baseado nas especificações
- **Uso**: `/t2c.implement specs/001-exemplo`

### `/t2c.validate [caminho]`
- **Arquivo de referência**: `.vscode/commands/t2c.validate.md`
- **Função**: Valida estrutura e completude dos arquivos de especificação
- **Uso**: `/t2c.validate specs/001-exemplo`

## Como Funcionar

Quando o usuário usar um slash command:

1. **Detecte o comando**: Se começar com `/t2c.`, é um comando customizado
2. **Leia o arquivo**: Abra e leia o conteúdo de `.vscode/commands/[comando].md`
3. **Siga as instruções**: Execute EXATAMENTE o que está descrito no arquivo
4. **Respeite as regras**: Preste atenção especial às seções "⚠️ REGRA ABSOLUTA"

## Importante

- **NUNCA crie scripts Python** quando o comando pedir para executar um script existente
- **SEMPRE leia o arquivo markdown** antes de executar qualquer ação
- **Siga as instruções passo a passo** conforme descrito nos arquivos de comando
- **Use os templates** em `.specify/templates/` como referência quando necessário

## Estrutura de Arquivos

```
.vscode/
└── commands/
    ├── t2c.extract-ddp.md  # Instruções completas para extrair DDP
    ├── t2c.tasks.md         # Instruções para gerar tasks.md
    ├── t2c.implement.md     # Instruções para implementar framework
    └── t2c.validate.md      # Instruções para validar specs
```

Cada arquivo contém instruções detalhadas sobre como executar o comando correspondente.
"""
    
    (vscode_dir / "copilot-instructions.md").write_text(instructions_content, encoding="utf-8")


def _create_vscode_tasks(vscode_dir: Path):
    """Cria tasks.json com tasks para executar os scripts"""
    tasks = {
        "version": "2.0.0",
        "tasks": [
            {
                "label": "T2C: Extract DDP",
                "type": "shell",
                "command": "python",
                "args": [
                    "${workspaceFolder}/.specify/scripts/extract-ddp.py"
                ],
                "problemMatcher": [],
                "presentation": {
                    "reveal": "always",
                    "panel": "new"
                },
                "group": {
                    "kind": "build",
                    "isDefault": False
                }
            },
            {
                "label": "T2C: Extract DDP (with file)",
                "type": "shell",
                "command": "python",
                "args": [
                    "${workspaceFolder}/.specify/scripts/extract-ddp.py",
                    "${input:ddpPath}"
                ],
                "problemMatcher": [],
                "presentation": {
                    "reveal": "always",
                    "panel": "new"
                },
                "group": {
                    "kind": "build",
                    "isDefault": False
                }
            }
        ],
        "inputs": [
            {
                "id": "ddpPath",
                "type": "promptString",
                "description": "Caminho do arquivo DDP (relativo ao workspace)",
                "default": "DDP/ddp.pptx"
            }
        ]
    }
    
    import json
    (vscode_dir / "tasks.json").write_text(
        json.dumps(tasks, indent=2, ensure_ascii=False),
        encoding="utf-8"
    )


def _create_vscode_readme(vscode_dir: Path):
    """Cria README explicando como usar os comandos com GitHub Copilot"""
    readme_content = """# Comandos T2C para VS Code + GitHub Copilot

Este diretório contém os comandos T2C disponíveis para uso com GitHub Copilot.

## Como Usar

### Método 1: Slash Commands (Igual ao Cursor) ⭐

No chat do GitHub Copilot, use os slash commands diretamente:

- **Extrair DDP**: `/t2c.extract-ddp` ou `/t2c.extract-ddp specs/001-exemplo/DDP/ddp.pptx`
- **Gerar Tasks**: `/t2c.tasks specs/001-exemplo`
- **Implementar Framework**: `/t2c.implement specs/001-exemplo`
- **Validar Specs**: `/t2c.validate specs/001-exemplo`

O Copilot reconhecerá os slash commands e lerá automaticamente os arquivos em `.vscode/commands/` para entender o que fazer.

**Nota**: O arquivo `.vscode/copilot-instructions.md` contém instruções para o Copilot sobre como processar esses comandos.

### Método 2: Mencionar ao GitHub Copilot

Você também pode mencionar o comando diretamente:

- **Extrair DDP**: "Execute o comando t2c.extract-ddp" ou "Extrair DDP usando t2c.extract-ddp"
- **Gerar Tasks**: "Execute o comando t2c.tasks" ou "Gerar tasks usando t2c.tasks"
- **Implementar Framework**: "Execute o comando t2c.implement" ou "Implementar framework usando t2c.implement"
- **Validar Specs**: "Execute o comando t2c.validate" ou "Validar specs usando t2c.validate"

### Método 3: Usar Tasks do VS Code

1. Pressione `Ctrl+Shift+P` (ou `Cmd+Shift+P` no Mac)
2. Digite "Tasks: Run Task"
3. Selecione uma das tasks disponíveis:
   - **T2C: Extract DDP** - Extrai DDP automaticamente
   - **T2C: Extract DDP (with file)** - Extrai DDP de um arquivo específico

### Método 4: Executar Scripts Diretamente

Você também pode executar os scripts diretamente no terminal:

```bash
# Extrair DDP (procura automaticamente)
python .specify/scripts/extract-ddp.py

# Extrair DDP de arquivo específico
python .specify/scripts/extract-ddp.py DDP/arquivo.pptx
```

## Comandos Disponíveis

### t2c.extract-ddp

Extrai o texto de todos os slides de um arquivo DDP.pptx.

**Uso com Copilot:**
- "Execute t2c.extract-ddp"
- "Extrair DDP do arquivo specs/001-exemplo/DDP/ddp.pptx"

### t2c.tasks

Gera o arquivo tasks.md baseado em spec.md e business-rules.md.

**Uso com Copilot:**
- "Execute t2c.tasks para specs/001-exemplo"
- "Gerar tasks.md baseado nas specs"

### t2c.implement

Gera o framework T2C completo baseado nas especificações.

**Uso com Copilot:**
- "Execute t2c.implement para specs/001-exemplo"
- "Implementar framework T2C completo"

### t2c.validate

Valida a estrutura e completude dos arquivos de especificação.

**Uso com Copilot:**
- "Execute t2c.validate para specs/001-exemplo"
- "Validar todas as specs"

## Documentação Completa

Consulte os arquivos em `.vscode/commands/` para documentação detalhada de cada comando:
- `t2c.extract-ddp.md`
- `t2c.tasks.md`
- `t2c.implement.md`
- `t2c.validate.md`

## Dicas

1. **Sempre mencione o comando completo**: "t2c.extract-ddp" em vez de apenas "extrair"
2. **Seja específico sobre o caminho**: "t2c.extract-ddp specs/001-exemplo/DDP/ddp.pptx"
3. **Leia a documentação**: O Copilot pode ler os arquivos `.md` para entender melhor o que fazer
4. **Use as tasks**: Para execução rápida, use as tasks do VS Code (`Ctrl+Shift+P` > "Tasks: Run Task")
"""
    
    (vscode_dir / "README.md").write_text(readme_content, encoding="utf-8")


def _create_initial_files(project_path: Path, project_name: str):
    """Cria arquivos iniciais do projeto"""
    # README.md
    readme_content = f"""# {project_name}

Projeto de automação RPA criado com RPA Spec-Kit.

## Estrutura do Projeto

\`\`\`
{project_name}/
├── .specify/          # Configurações e templates
│   ├── memory/        # Constitution do framework T2C
│   ├── templates/     # Templates de especificação
│   └── scripts/       # Scripts prontos (extract-ddp.py)
├── specs/             # Especificações de automações
│   └── 001-[nome]/    # Primeira automação
│       ├── spec.md     # ARQUIVO PRINCIPAL - Arquitetura completa
│       ├── tests.md    # Cenários de teste e validações
│       ├── selectors.md
│       ├── business-rules.md
│       ├── tasks.md
│       └── DDP/        # DDPs (Documentos de Design de Processo)
├── generated/         # Framework T2C gerado
└── DDP/               # DDPs gerais
\`\`\`

## Fluxo de Trabalho

1. **Inicialização**: Projeto já inicializado ✓
2. **Extrair DDP**: Coloque DDP.pptx em `specs/001-[nome]/DDP/` ou `DDP/` e execute `/t2c.extract-ddp`
   - O script instala dependências automaticamente se necessário
4. **Completar Specs**: Revise e complete os arquivos .md gerados
5. **Gerar Tasks** (Opcional): Execute `/t2c.tasks` para gerar tasks.md
6. **Implementar**: Execute `/t2c.implement` para gerar o framework T2C completo

## Comandos Disponíveis

- `/t2c.extract-ddp` - Extrai informações de DDP.pptx
- `/t2c.tasks` - Gera tasks.md baseado nas specs
- `/t2c.implement` - Gera framework T2C completo
- `/t2c.validate` - Valida estrutura e completude das specs

## Próximos Passos

1. Crie uma nova feature: `specs/001-[nome-da-automacao]/`
2. Coloque o DDP.pptx na pasta DDP/
3. Execute `/t2c.extract-ddp` para extrair informações
4. Complete os arquivos .md conforme necessário
5. Execute `/t2c.implement` para gerar o framework
"""
    (project_path / "README.md").write_text(readme_content, encoding="utf-8")
    
    # .gitignore
    gitignore_content = """# Python
__pycache__/
*.py[cod]
*$py.class
*.so
.Python
build/
develop-eggs/
dist/
downloads/
eggs/
.eggs/
lib/
lib64/
parts/
sdist/
var/
wheels/
*.egg-info/
.installed.cfg
*.egg

# Virtual Environment
venv/
env/
ENV/
.venv

# IDE
.vscode/
.idea/
*.swp
*.swo
*~

# RPA Spec-Kit
generated/
*.pptx
*.xlsx
*.db
*.sqlite

# Logs
*.log

# OS
.DS_Store
Thumbs.db
"""
    (project_path / ".gitignore").write_text(gitignore_content, encoding="utf-8")

