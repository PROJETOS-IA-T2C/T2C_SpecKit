"""
Extrator de DDP - Extrai informações de arquivos PPTX
"""
from pathlib import Path
from pptx import Presentation
from typing import Dict, List, Optional
import re


class DDPExtractor:
    """Classe para extrair informações de DDP.pptx"""
    
    def __init__(self, pptx_path: str):
        """
        Inicializa o extrator
        
        Args:
            pptx_path: Caminho para o arquivo DDP.pptx
        """
        self.pptx_path = Path(pptx_path)
        if not self.pptx_path.exists():
            raise FileNotFoundError(f"DDP não encontrado: {pptx_path}")
        
        self.presentation = Presentation(str(self.pptx_path))
        self.slides_text: List[str] = []
        self.spec_info: Dict = {}
        self.selectors: List[Dict] = []
        self.business_rules: List[Dict] = []
    
    def extract_from_pptx(self) -> Dict:
        """
        Extrai todas as informações do PPTX
        
        Returns:
            Dicionário com todas as informações extraídas
        """
        self.extract_text_from_slides()
        self.parse_spec_info()
        self.parse_selectors()
        self.parse_business_rules()
        
        return {
            'spec_info': self.spec_info,
            'selectors': self.selectors,
            'business_rules': self.business_rules,
            'raw_text': '\n\n'.join(self.slides_text)
        }
    
    def extract_text_from_slides(self):
        """Extrai texto de todos os slides"""
        for slide in self.presentation.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            self.slides_text.append('\n'.join(slide_text))
    
    def parse_spec_info(self):
        """Identifica informações de spec do DDP"""
        full_text = '\n\n'.join(self.slides_text).lower()
        
        # Tentar identificar cenários de usuário
        scenarios = re.findall(r'(?:cenário|scenario|como|given|when|then).*?(?=\n\n|\Z)', full_text, re.IGNORECASE | re.DOTALL)
        
        # Tentar identificar requisitos
        requirements = re.findall(r'(?:requisito|requirement|rf\d+|rnf\d+).*?(?=\n\n|\Z)', full_text, re.IGNORECASE | re.DOTALL)
        
        # Tentar identificar critérios de sucesso
        success_criteria = re.findall(r'(?:critério|criteria|sucesso|success).*?(?=\n\n|\Z)', full_text, re.IGNORECASE | re.DOTALL)
        
        self.spec_info = {
            'scenarios': scenarios[:5] if scenarios else [],
            'requirements': requirements[:10] if requirements else [],
            'success_criteria': success_criteria[:10] if success_criteria else [],
            'raw_text': '\n\n'.join(self.slides_text)
        }
    
    def parse_selectors(self):
        """Identifica seletores mencionados no DDP"""
        full_text = '\n\n'.join(self.slides_text)
        
        # Procurar por padrões de seletores
        selector_patterns = [
            r'(?:seletor|selector|locator|elemento|element).*?(?=\n\n|\Z)',
            r'(?:botão|button|btn).*?(?=\n\n|\Z)',
            r'(?:campo|field|input).*?(?=\n\n|\Z)',
            r'(?:tabela|table).*?(?=\n\n|\Z)',
        ]
        
        for pattern in selector_patterns:
            matches = re.findall(pattern, full_text, re.IGNORECASE | re.DOTALL)
            for match in matches[:5]:  # Limitar a 5 por padrão
                self.selectors.append({
                    'description': match.strip(),
                    'type': 'unknown'
                })
    
    def parse_business_rules(self):
        """Identifica regras de negócio no DDP"""
        full_text = '\n\n'.join(self.slides_text)
        
        # Procurar por padrões de regras
        rule_patterns = [
            r'(?:val\d+|validação|validation).*?(?=\n\n|\Z)',
            r'(?:cond\d+|condição|condition).*?(?=\n\n|\Z)',
            r'(?:reg\d+|regra|rule).*?(?=\n\n|\Z)',
            r'(?:se.*?então|if.*?then).*?(?=\n\n|\Z)',
        ]
        
        for pattern in rule_patterns:
            matches = re.findall(pattern, full_text, re.IGNORECASE | re.DOTALL)
            for match in matches[:10]:  # Limitar a 10 por padrão
                rule_type = 'validation' if 'val' in match.lower() else \
                           'condition' if 'cond' in match.lower() else \
                           'rule' if 'reg' in match.lower() else 'unknown'
                
                self.business_rules.append({
                    'description': match.strip(),
                    'type': rule_type
                })
    
    def fill_spec_md(self, output_path: Path) -> str:
        """
        Preenche spec.md com informações extraídas
        
        Args:
            output_path: Caminho onde salvar spec.md
            
        Returns:
            Conteúdo do spec.md gerado
        """
        content = "# Especificação de Automação RPA\n\n"
        content += "> **Nota:** Este arquivo foi gerado automaticamente a partir do DDP. Revise e complete conforme necessário.\n\n"
        
        content += "## User Scenarios\n\n"
        if self.spec_info.get('scenarios'):
            for i, scenario in enumerate(self.spec_info['scenarios'], 1):
                content += f"### Scenario {i}: [AUTO] Extraído do DDP\n"
                content += f"{scenario}\n\n"
        else:
            content += "### Scenario 1: [Preencher manualmente]\n"
            content += "**Como** [persona]\n"
            content += "**Eu quero** [ação]\n"
            content += "**Para que** [objetivo]\n\n"
        
        content += "## Requirements\n\n"
        content += "### Funcionais\n"
        if self.spec_info.get('requirements'):
            for req in self.spec_info['requirements'][:10]:
                content += f"- [ ] [AUTO] {req[:100]}...\n"
        else:
            content += "- [ ] RF001: [Descrição]\n"
        
        content += "\n### Não Funcionais\n"
        content += "- [ ] RNF001: [Descrição]\n\n"
        
        content += "## Success Criteria\n\n"
        if self.spec_info.get('success_criteria'):
            for i, criteria in enumerate(self.spec_info['success_criteria'], 1):
                content += f"- [ ] [AUTO] SC{i:03d}: {criteria[:100]}...\n"
        else:
            content += "- [ ] SC001: [Critério de sucesso]\n\n"
        
        content += "## Key Entities\n\n"
        content += "### Entidade 1: [Nome]\n"
        content += "- Campo1: [Tipo/Descrição]\n"
        content += "- Campo2: [Tipo/Descrição]\n\n"
        
        content += "## Observações\n\n"
        content += "[Observações adicionais]\n"
        
        return content
    
    def fill_plan_md(self, output_path: Path) -> str:
        """
        Preenche plan.md com stack padrão
        
        Args:
            output_path: Caminho onde salvar plan.md
            
        Returns:
            Conteúdo do plan.md gerado
        """
        content = "# Plano Técnico de Implementação\n\n"
        content += "> **Nota:** Este arquivo foi gerado automaticamente. Revise e complete conforme necessário.\n\n"
        
        content += "## Stack Tecnológica\n\n"
        content += "- **Framework:** T2C Framework (v2.2.3)\n"
        content += "- **Automação Web:** Clicknium\n"
        content += "- **Plataforma:** BotCity\n"
        content += "- **Linguagem:** Python 3.8+\n\n"
        
        content += "## Arquitetura do Robô\n\n"
        content += "### Componentes Principais\n\n"
        content += "1. **T2CProcess**\n"
        content += "   - Responsabilidade: [Descrição]\n"
        content += "   - Fluxo: [Descrição do fluxo]\n\n"
        content += "2. **T2CInitAllApplications**\n"
        content += "   - Responsabilidade: [Descrição]\n"
        content += "   - Aplicações a inicializar: [Lista]\n\n"
        content += "3. **T2CCloseAllApplications**\n"
        content += "   - Responsabilidade: [Descrição]\n"
        content += "   - Aplicações a fechar: [Lista]\n\n"
        
        content += "## Integrações\n\n"
        content += "- [ ] Maestro (BotCity)\n"
        content += "- [ ] T2CTracker\n"
        content += "- [ ] Clicknium\n"
        content += "- [ ] E-mail\n\n"
        
        content += "## Estrutura de Dados\n\n"
        content += "### Fila de Processamento\n"
        content += "- Referência: [Campo identificador]\n"
        content += "- Info Adicionais: [Estrutura JSON]\n\n"
        
        content += "## Fluxo de Execução\n\n"
        content += "1. [Passo 1]\n"
        content += "2. [Passo 2]\n"
        content += "3. [Passo 3]\n"
        
        return content
    
    def fill_selectors_md(self, output_path: Path) -> str:
        """
        Preenche selectors.md com seletores identificados
        
        Args:
            output_path: Caminho onde salvar selectors.md
            
        Returns:
            Conteúdo do selectors.md gerado
        """
        content = "# Seletores Clicknium\n\n"
        content += "> **Nota:** Este arquivo foi gerado automaticamente. Revise e complete conforme necessário.\n\n"
        
        content += "## Estrutura de Locators\n\n"
        
        if self.selectors:
            content += "### Pasta: [nome_pasta] [AUTO]\n\n"
            for i, selector in enumerate(self.selectors[:10], 1):
                content += f"#### elemento_{i} [AUTO]\n"
                content += f"- **Tipo:** [button/input/div/etc]\n"
                content += f"- **Seletor:** {selector['description'][:100]}...\n"
                content += f"- **Uso:** [onde é usado]\n\n"
        else:
            content += "### Pasta: [nome_pasta]\n\n"
            content += "#### [nome_elemento]\n"
            content += "- **Tipo:** [button/input/div/etc]\n"
            content += "- **Seletor:** [descrição do seletor]\n"
            content += "- **Uso:** [onde é usado]\n\n"
        
        content += "## Notas\n\n"
        content += "- Todos os seletores devem ser criados no Clicknium Recorder\n"
        content += "- Manter nomenclatura consistente\n"
        content += "- Documentar mudanças de UI que afetam seletores\n"
        
        return content
    
    def fill_business_rules_md(self, output_path: Path) -> str:
        """
        Preenche business-rules.md com regras identificadas
        
        Args:
            output_path: Caminho onde salvar business-rules.md
            
        Returns:
            Conteúdo do business-rules.md gerado
        """
        content = "# Regras de Negócio\n\n"
        content += "> **Nota:** Este arquivo foi gerado automaticamente. Revise e complete conforme necessário.\n\n"
        
        # Separar por tipo
        validations = [r for r in self.business_rules if r['type'] == 'validation']
        conditions = [r for r in self.business_rules if r['type'] == 'condition']
        rules = [r for r in self.business_rules if r['type'] == 'rule']
        
        content += "## Validações (VAL*)\n\n"
        if validations:
            for i, val in enumerate(validations[:10], 1):
                content += f"### VAL{i:03d}: [AUTO] Extraído do DDP\n"
                content += f"- **Descrição:** {val['description'][:200]}...\n"
                content += f"- **Condição:** [Quando aplicar]\n"
                content += f"- **Ação em Erro:** [O que fazer se falhar]\n"
                content += f"- **Exceção:** BusinessRuleException\n\n"
        else:
            content += "### VAL001: [Nome da Validação]\n"
            content += "- **Descrição:** [Descrição completa]\n"
            content += "- **Condição:** [Quando aplicar]\n"
            content += "- **Ação em Erro:** [O que fazer se falhar]\n"
            content += "- **Exceção:** BusinessRuleException\n\n"
        
        content += "## Condições Especiais (COND*)\n\n"
        if conditions:
            for i, cond in enumerate(conditions[:10], 1):
                content += f"### COND{i:03d}: [AUTO] Extraído do DDP\n"
                content += f"- **Descrição:** {cond['description'][:200]}...\n"
                content += f"- **Condição:** [Quando aplicar]\n"
                content += f"- **Ação:** [O que fazer]\n\n"
        else:
            content += "### COND001: [Nome da Condição]\n"
            content += "- **Descrição:** [Descrição completa]\n"
            content += "- **Condição:** [Quando aplicar]\n"
            content += "- **Ação:** [O que fazer]\n\n"
        
        content += "## Regras de Processamento (REG*)\n\n"
        if rules:
            for i, reg in enumerate(rules[:10], 1):
                content += f"### REG{i:03d}: [AUTO] Extraído do DDP\n"
                content += f"- **Descrição:** {reg['description'][:200]}...\n"
                content += f"- **Aplicação:** [Quando aplicar]\n"
                content += f"- **Resultado:** [Resultado esperado]\n\n"
        else:
            content += "### REG001: [Nome da Regra]\n"
            content += "- **Descrição:** [Descrição completa]\n"
            content += "- **Aplicação:** [Quando aplicar]\n"
            content += "- **Resultado:** [Resultado esperado]\n\n"
        
        return content


def extract_ddp(pptx_path: str) -> str:
    """
    Função principal para extrair texto do DDP
    
    Args:
        pptx_path: Caminho para o DDP.pptx
        
    Returns:
        Texto extraído de todos os slides formatado para a LLM
    """
    extractor = DDPExtractor(pptx_path)
    extractor.extract_text_from_slides()
    
    # Formatar texto para apresentar à LLM
    formatted_text = "# Conteúdo Extraído do DDP\n\n"
    formatted_text += f"**Arquivo:** {pptx_path}\n\n"
    formatted_text += f"**Total de slides:** {len(extractor.slides_text)}\n\n"
    formatted_text += "---\n\n"
    
    for i, slide_text in enumerate(extractor.slides_text, 1):
        formatted_text += f"## Slide {i}\n\n"
        formatted_text += f"{slide_text}\n\n"
        formatted_text += "---\n\n"
    
    return formatted_text

